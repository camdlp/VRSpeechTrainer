from flask import Flask, jsonify, request, send_file, render_template
import os
import io
from pptx import Presentation
from PyPDF2 import PdfFileReader
from PIL import Image

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads/'

###### CONVERSIONS ######

def convert_pptx_to_png(filepath):
    """
    Convierte un archivo .pptx a una imagen PNG
    """
    # Cargar presentación
    prs = Presentation(filepath)

    # Obtener la primera diapositiva
    slide = prs.slides[0]

    # Generar imagen
    img_bytes = io.BytesIO()
    slide.shapes[0].image.save(img_bytes, "png")
    img_bytes.seek(0)
    img = Image.open(img_bytes)

    return img


def convert_pdf_to_png(filepath):
    """
    Convierte un archivo .pdf a una imagen PNG
    """
    # Cargar PDF
    with open(filepath, "rb") as f:
        pdf = PdfFileReader(f)

        # Obtener la primera página
        page = pdf.getPage(0)

        # Generar imagen
        img_bytes = io.BytesIO()
        page.thumbnail((1000, 1000))
        page.compressContentStreams()
        page.writeToStream(img_bytes)
        img_bytes.seek(0)
        img = Image.open(img_bytes)

    return img


def convert_to_png(filepath):
    """
    Convierte un archivo a una imagen PNG
    """
    # Obtener la extensión del archivo
    ext = os.path.splitext(filepath)[1].lower()

    # Realizar la conversión correspondiente
    if ext == ".pptx":
        return convert_pptx_to_png(filepath)
    elif ext == ".pdf":
        return convert_pdf_to_png(filepath)
    else:
        raise ValueError("Archivo no compatible")


###### ROUTES ######

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/slides')
def get_slide_files():
    folder_path = './Slides'
    files = os.listdir(folder_path)
    return jsonify(files)



@app.route('/recordings')
def get_recording_files():
    folder_path = './Recordings'
    files = os.listdir(folder_path)
    return jsonify(files)


@app.route('/api/download')
def download_file():
    filename = request.args.get('filename')
    filepath = os.path.join(request.args.get('type'), filename)
    # Log the file path
    print(filepath)
    if os.path.isfile(filepath):
        return send_file(filepath, as_attachment=True)
    else:
        return 'File not found', 404

@app.route('/api/upload', methods=['POST'])
def upload_file():
    file = request.files['file']
    if file:
        # Obtén la extensión del archivo
        file_ext = os.path.splitext(file.filename)[1].lower()

        # Determina la carpeta de destino según la extensión del archivo
        if file_ext in ('.pptx', '.pdf'):
            type = 'Slides'
        elif file_ext == '.mp4':
            type = 'Recordings'
        else:
            return 'File type not supported', 400

        upload_dir = type
        if not os.path.isdir(upload_dir):
            os.makedirs(upload_dir)
        file.save(os.path.join(upload_dir, file.filename))
        return 'File uploaded successfully'
    else:
        return 'No file uploaded', 400
    
if __name__ == '__main__':
    app.run(debug=True)
